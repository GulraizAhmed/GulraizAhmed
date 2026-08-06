#!/usr/bin/env python3
"""Generate Punjab Zameen Application PowerPoint presentation."""

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


# PLRA / Punjab Zameen brand colors (from attached designs)
GREEN = RGBColor(0x00, 0x68, 0x37)
GREEN_DARK = RGBColor(0x0F, 0x4A, 0x32)
GREEN_MID = RGBColor(0x1B, 0x7A, 0x4A)
GREEN_SOFT = RGBColor(0xE8, 0xF5, 0xEE)
GREEN_BORDER = RGBColor(0x2F, 0x9E, 0x6B)
TITLE_NAVY = RGBColor(0x1A, 0x2B, 0x4A)
GRAY = RGBColor(0x55, 0x55, 0x55)
GRAY_LIGHT = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
NSMAP_P = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}


def set_run(run, size=14, bold=False, color=BLACK, italic=False, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_text(shape, lines, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(valign, "t"))
    except Exception:
        pass
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, size, bold, color, italic = item, 14, False, BLACK, False
        else:
            text = item[0]
            size = item[1]
            bold = item[2]
            color = item[3]
            italic = item[4] if len(item) > 4 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = text
        set_run(run, size=size, bold=bold, color=color, italic=italic)


def add_footer(slide, prs):
    box = slide.shapes.add_textbox(
        Inches(0.5),
        prs.slide_height - Inches(0.42),
        prs.slide_width - Inches(1.0),
        Inches(0.28),
    )
    add_text(
        box,
        [("PUNJAB LAND RECORDS AUTHORITY (PLRA)  |  GOVT OF THE PUNJAB", 10, False, GRAY_LIGHT)],
        align=PP_ALIGN.CENTER,
    )


def add_eyebrow(slide, text="CORE APPLICATION MODULES"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(8), Inches(0.28))
    add_text(box, [(text, 11, True, GREEN)])


def add_title_block(slide, title, subtitle=None, y=0.48):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.3), Inches(0.48))
    add_text(box, [(title, 26, True, TITLE_NAVY)])
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.42), Inches(12.3), Inches(0.32))
        add_text(box2, [(subtitle, 13, False, GRAY)])


def add_rounded_rect(slide, left, top, width, height, fill=WHITE, line=GREEN_BORDER, line_width=1.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_circle(slide, left, top, size, fill=GREEN, text=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if text is not None:
        add_text(shape, [(str(text), 13, True, WHITE)], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        for p in shape.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
    return shape


def add_bullet_card(slide, left, top, width, height, title, bullets):
    add_rounded_rect(slide, left, top, width, height, fill=WHITE, line=GREEN_BORDER)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

    t = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.14), width - Inches(0.4), Inches(0.34))
    add_text(t, [(title, 14, True, GREEN_DARK)])

    body = slide.shapes.add_textbox(
        left + Inches(0.25), top + Inches(0.48), width - Inches(0.4), height - Inches(0.62)
    )
    lines = [(f"•  {b}", 12, False, GRAY) for b in bullets]
    add_text(body, lines)
    return body


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_transition(slide, kind="fade"):
    """Add a formal slide transition (Fade / Push / Wipe)."""
    sld = slide._element
    for child in list(sld):
        if child.tag == qn("p:transition"):
            sld.remove(child)

    transition = etree.Element(qn("p:transition"))
    transition.set("spd", "med")
    transition.set("advClick", "1")

    if kind == "push":
        node = etree.SubElement(transition, qn("p:push"))
        node.set("dir", "l")
    elif kind == "wipe":
        node = etree.SubElement(transition, qn("p:wipe"))
        node.set("dir", "l")
    else:
        etree.SubElement(transition, qn("p:fade"))

    # Insert transition before clrMapOvr if present, else append
    sld.append(transition)


def add_appear_animation(slide, shape, order):
    """Add a subtle Appear animation on click for a shape."""
    # Build / reuse timing tree
    cSld = slide._element.find(qn("p:cSld"))
    # timing lives on sld, after cSld
    sld = slide._element
    timing = sld.find(qn("p:timing"))
    if timing is None:
        timing = etree.Element(qn("p:timing"))
        # place after transition if any, else after cSld
        sld.append(timing)
        tn_list = etree.SubElement(timing, qn("p:tnLst"))
        par = etree.SubElement(tn_list, qn("p:par"))
        ctn = etree.SubElement(par, qn("p:cTn"))
        ctn.set("id", "1")
        ctn.set("dur", "indefinite")
        ctn.set("restart", "never")
        ctn.set("nodeType", "tmRoot")
        child_tn = etree.SubElement(ctn, qn("p:childTnLst"))
        seq = etree.SubElement(child_tn, qn("p:seq"))
        seq.set("concurrent", "1")
        seq.set("nextAc", "seek")
        seq_ctn = etree.SubElement(seq, qn("p:cTn"))
        seq_ctn.set("id", "2")
        seq_ctn.set("dur", "indefinite")
        seq_ctn.set("nodeType", "mainSeq")
        etree.SubElement(seq_ctn, qn("p:childTnLst"))
        prev = etree.SubElement(seq, qn("p:prevCondLst"))
        cond = etree.SubElement(prev, qn("p:cond"))
        cond.set("evt", "onPrev")
        cond.set("delay", "0")
        tgt = etree.SubElement(cond, qn("p:tgtEl"))
        etree.SubElement(tgt, qn("p:sldTgt"))
        nxt = etree.SubElement(seq, qn("p:nextCondLst"))
        cond2 = etree.SubElement(nxt, qn("p:cond"))
        cond2.set("evt", "onNext")
        cond2.set("delay", "0")
        tgt2 = etree.SubElement(cond2, qn("p:tgtEl"))
        etree.SubElement(tgt2, qn("p:sldTgt"))

    # Find mainSeq childTnLst
    child_lists = timing.findall(".//{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst")
    if not child_lists:
        return
    # Prefer the deepest childTnLst under mainSeq
    main_list = None
    for cl in child_lists:
        parent = cl.getparent()
        if parent is not None and parent.get("nodeType") == "mainSeq":
            main_list = cl
            break
    if main_list is None:
        main_list = child_lists[-1]

    anim_id = 10 + order * 2
    spid = str(shape.shape_id)

    par = etree.SubElement(main_list, qn("p:par"))
    ctn = etree.SubElement(par, qn("p:cTn"))
    ctn.set("id", str(anim_id))
    ctn.set("fill", "hold")
    if order == 1:
        ctn.set("delay", "0")
    else:
        ctn.set("delay", "0")
    ctn.set("grpId", "0")
    ctn.set("nodeType", "clickEffect")
    st = etree.SubElement(ctn, qn("p:stCondLst"))
    cond = etree.SubElement(st, qn("p:cond"))
    cond.set("delay", "0")
    child = etree.SubElement(ctn, qn("p:childTnLst"))

    anim = etree.SubElement(child, qn("p:animEffect"))
    anim.set("transition", "in")
    anim.set("filter", "fade")
    cBhvr = etree.SubElement(anim, qn("p:cBhvr"))
    cTn2 = etree.SubElement(cBhvr, qn("p:cTn"))
    cTn2.set("id", str(anim_id + 1))
    cTn2.set("dur", "500")
    tgt = etree.SubElement(cBhvr, qn("p:tgtEl"))
    sp = etree.SubElement(tgt, qn("p:spTgt"))
    sp.set("spid", spid)


def apply_transitions(prs):
    kinds = ["fade", "push", "wipe"]
    for i, slide in enumerate(prs.slides):
        set_transition(slide, kinds[i % len(kinds)])


# ---------- Slides ----------

def slide_title(prs):
    s = blank_slide(prs)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(3.15))
    band.fill.solid()
    band.fill.fore_color.rgb = GREEN
    band.line.fill.background()

    soft = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.15), prs.slide_width, Inches(0.12))
    soft.fill.solid()
    soft.fill.fore_color.rgb = GREEN_MID
    soft.line.fill.background()

    box = s.shapes.add_textbox(Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.35))
    add_text(box, [("PUNJAB LAND RECORDS AUTHORITY (PLRA)", 13, True, WHITE)], align=PP_ALIGN.CENTER)

    box = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.7))
    add_text(box, [("Punjab Zameen Application", 38, True, WHITE)], align=PP_ALIGN.CENTER)

    box = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.4))
    add_text(
        box,
        [("Citizen-Facing Digital Land Services Platform", 17, False, RGBColor(0xD8, 0xF0, 0xE4))],
        align=PP_ALIGN.CENTER,
    )

    points = [
        ("11+ Integrated Systems", "Enterprise interoperability across provincial departments"),
        ("End-to-End Digital Services", "GPC, Record Copy, Payments, Complaints & more"),
        ("Citizen Empowerment", "Verified holdings, online payments, and status tracking"),
    ]
    x = Inches(0.7)
    for title, desc in points:
        add_rounded_rect(s, x, Inches(3.7), Inches(3.9), Inches(1.7), fill=GREEN_SOFT, line=GREEN_BORDER)
        t = s.shapes.add_textbox(x + Inches(0.2), Inches(3.9), Inches(3.5), Inches(0.4))
        add_text(t, [(title, 14, True, GREEN_DARK)])
        d = s.shapes.add_textbox(x + Inches(0.2), Inches(4.4), Inches(3.5), Inches(0.75))
        add_text(d, [(desc, 12, False, GRAY)])
        x += Inches(4.1)

    add_footer(s, prs)


def slide_agenda(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "OVERVIEW")
    add_title_block(s, "Agenda", "Formal walkthrough of Punjab Zameen modules and workflows")

    left_items = [
        "1. Why Punjab Zameen Remodeling?",
        "2. Integrated Enterprise Ecosystem",
        "3. Home Screen Service Cards",
        "4. My Properties & GPC Workflow",
        "5. Approved Housing Society (HSMS)",
        "6. Get Record Copy",
        "7. Payment Gateway",
        "8. Complaints Module (CMS)",
        "9. Arazi Muawin",
    ]
    right_items = [
        "10. Revenue Court Cases (RCMS)",
        "11. e-Registration",
        "12. e-Stamp",
        "13. Property Tax (PTS)",
        "14. e-Leasing",
        "15. Community Module",
        "16. Digital Vault",
        "17. Notifications & Messaging",
        "18. Closing",
    ]

    for i, items in enumerate([left_items, right_items]):
        x = Inches(0.55) if i == 0 else Inches(6.85)
        add_rounded_rect(s, x, Inches(1.45), Inches(5.9), Inches(5.15), fill=WHITE, line=GREEN_BORDER)
        box = s.shapes.add_textbox(x + Inches(0.35), Inches(1.65), Inches(5.2), Inches(4.7))
        add_text(box, [(item, 14, False, BLACK) for item in items])

    add_footer(s, prs)


def slide_why_remodel(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "STRATEGIC RATIONALE")
    add_title_block(s, "Why Punjab Zameen Remodeling?", "From fragmented legacy access to a unified citizen platform")

    cards = [
        ("Citizen Empowerment", [
            "Direct access to verified property holdings",
            "GPC requests and tax payments without middlemen",
            "Self-service digital land journeys",
        ]),
        ("Enterprise Interoperability", [
            "Standardized RESTful APIs",
            "Connecting 11+ provincial government systems",
            "Single front-end for PLRA services",
        ]),
        ("Financial Transparency", [
            "Automated fee calculation",
            "PSID tracking and reconciliation",
            "Duplicate-payment safeguards",
        ]),
        ("Access to PLRA Services", [
            "Multiple PLRA services from one platform",
            "Reduced citizen load on offices",
            "Live status, vault storage, and guidance",
        ]),
    ]
    positions = [
        (Inches(0.5), Inches(1.45)),
        (Inches(6.7), Inches(1.45)),
        (Inches(0.5), Inches(4.0)),
        (Inches(6.7), Inches(4.0)),
    ]
    shapes = []
    for (x, y), (title, bullets) in zip(positions, cards):
        add_bullet_card(s, x, y, Inches(5.9), Inches(2.25), title, bullets)
        shapes.append(s.shapes[-1])
    for i, sh in enumerate(shapes, 1):
        add_appear_animation(s, sh, i)
    add_footer(s, prs)


def slide_ecosystem(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "INTEGRATED ENTERPRISE ECOSYSTEM")
    add_title_block(
        s,
        "Seamless Interoperability Across 10+ Government Systems",
        "Punjab Zameen as the citizen front-end to PLRA and partner systems",
    )

    systems = [
        ("CIMS", "Central Identity\nManagement System"),
        ("CLRMIS", "Computerized Land Records\nManagement Information System"),
        ("Payment Gateway", "Online challan fee\ncollection & reconciliation"),
        ("HSMS", "Housing Society\nManagement System"),
        ("RCMS", "Revenue Court\nManagement System"),
        ("Arazi Muawin", "Franchise centers &\nland facilitation"),
        ("CMS", "Complaint\nManagement System"),
        ("e-Registration", "Property registration\npipeline & deed copies"),
        ("e-Stamping", "Stamp instrument\nfee & papers"),
        ("PTS", "Property Tax\nSystem"),
    ]

    start_x, start_y = Inches(0.45), Inches(1.5)
    card_w, card_h = Inches(2.4), Inches(1.55)
    gap_x, gap_y = Inches(0.2), Inches(0.18)

    for i, (code, desc) in enumerate(systems):
        row, col = divmod(i, 5)
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        add_rounded_rect(s, x, y, card_w, card_h, fill=GREEN_SOFT, line=GREEN_BORDER)
        t = s.shapes.add_textbox(x + Inches(0.1), y + Inches(0.22), card_w - Inches(0.2), Inches(0.4))
        add_text(t, [(code, 13, True, GREEN_DARK)], align=PP_ALIGN.CENTER)
        d = s.shapes.add_textbox(x + Inches(0.1), y + Inches(0.68), card_w - Inches(0.2), Inches(0.7))
        add_text(d, [(desc, 11, False, GRAY)], align=PP_ALIGN.CENTER)

    note = s.shapes.add_textbox(Inches(0.5), Inches(5.05), Inches(12.3), Inches(0.55))
    add_text(
        note,
        [
            (
                "Plus My Properties / GPC, Digital Vault, Community, and Notifications as platform capabilities "
                "connecting these systems for citizens.",
                12,
                True,
                GREEN,
            )
        ],
        align=PP_ALIGN.CENTER,
    )
    add_footer(s, prs)


def slide_home_cards(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "APPLICATION HOME SCREEN")
    add_title_block(s, "Service Cards on Punjab Zameen", "Each card opens a complete integrated system or module")

    cards = [
        ("My Properties", "GPC & actionable property hub"),
        ("Approved Housing Society", "HSMS — GIS & plot verification"),
        ("Get Record Copy", "Fard, Registry, Mutation, Girdawari"),
        ("Complaints", "CMS — register & track grievances"),
        ("Payments", "Payment Gateway for challans"),
        ("Revenue Court Cases", "RCMS — cases, cause lists, judgments"),
        ("Arazi Muawin", "Franchise locator & applications"),
        ("e-Registration", "Registration pipeline & deeds"),
        ("e-Stamp", "Stamp instruments & fees"),
        ("Property Tax", "PTS — tax search & settlement"),
        ("e-Leasing", "Schemes, deposits & live bidding"),
        ("Community", "Discussions, guides & practitioners"),
        ("Digital Vault", "Certificates, deeds & receipts"),
    ]

    start_x, start_y = Inches(0.4), Inches(1.4)
    card_w, card_h = Inches(3.1), Inches(0.95)
    gap_x, gap_y = Inches(0.15), Inches(0.12)

    for i, (title, desc) in enumerate(cards):
        row, col = divmod(i, 4)
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        if i == 12:
            x = start_x + 1.5 * (card_w + gap_x)
        add_rounded_rect(s, x, y, card_w, card_h, fill=WHITE, line=GREEN_BORDER)
        add_rounded_rect(
            s,
            x + Inches(0.12),
            y + Inches(0.2),
            Inches(0.55),
            Inches(0.55),
            fill=GREEN_SOFT,
            line=GREEN_SOFT,
        )
        t = s.shapes.add_textbox(x + Inches(0.8), y + Inches(0.14), card_w - Inches(0.95), Inches(0.35))
        add_text(t, [(title, 12, True, BLACK)])
        d = s.shapes.add_textbox(x + Inches(0.8), y + Inches(0.48), card_w - Inches(0.95), Inches(0.35))
        add_text(d, [(desc, 10, False, GRAY)])

    add_footer(s, prs)


def slide_my_properties(prs):
    s = blank_slide(prs)
    add_eyebrow(s)
    add_title_block(
        s,
        "My Properties Module",
        "Actionable property hub — from static list to service execution",
    )

    add_bullet_card(
        s,
        Inches(0.5),
        Inches(1.4),
        Inches(5.9),
        Inches(2.05),
        "AS-IS — Static List",
        [
            "Legacy app showed a plain static list of properties",
            "Zero functions or services could be performed",
            "No direct actions from property cards",
        ],
    )
    add_bullet_card(
        s,
        Inches(6.7),
        Inches(1.4),
        Inches(5.9),
        Inches(2.05),
        "TO-BE — Actionable Hub",
        [
            "Initiate Green Property Certificates (GPC)",
            "Request Fard / Registry copies",
            "Check court cases & transfer ownership",
        ],
    )

    tabs = s.shapes.add_textbox(Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.3))
    add_text(tabs, [("Property Categories (each tab shows total count)", 13, True, GREEN_DARK)])

    for i, (name, desc) in enumerate(
        [
            ("All", "Complete inventory of properties linked to the logged-in citizen"),
            ("Verified", "Properties confirmed against CLRMIS / ownership checks"),
            ("Unverified", "Properties awaiting verification or ownership confirmation"),
        ]
    ):
        x = Inches(0.5) + i * Inches(4.15)
        add_rounded_rect(s, x, Inches(4.05), Inches(3.95), Inches(1.75), fill=GREEN_SOFT, line=GREEN_BORDER)
        add_circle(s, x + Inches(0.2), Inches(4.25), Inches(0.38), text=str(i + 1))
        t = s.shapes.add_textbox(x + Inches(0.7), Inches(4.28), Inches(3.0), Inches(0.32))
        add_text(t, [(name, 15, True, GREEN_DARK)])
        d = s.shapes.add_textbox(x + Inches(0.2), Inches(4.75), Inches(3.55), Inches(0.85))
        add_text(d, [(desc, 12, False, GRAY)])

    add_footer(s, prs)


def slide_workflow(prs, eyebrow, title, subtitle, steps):
    """Horizontal numbered workflow cards matching attached GPC/CMS style."""
    s = blank_slide(prs)
    add_eyebrow(s, eyebrow)
    add_title_block(s, title, subtitle)

    n = len(steps)
    margin = Inches(0.32)
    gap = Inches(0.1)
    usable = prs.slide_width - 2 * margin - (n - 1) * gap
    card_w = usable / n
    card_h = Inches(4.4)
    top = Inches(1.4)

    for i, (step_title, desc, system) in enumerate(steps):
        x = margin + i * (card_w + gap)
        card = add_rounded_rect(s, x, top, card_w, card_h, fill=WHITE, line=GREEN_BORDER, line_width=1.75)
        circle_size = Inches(0.4)
        cx = x + (card_w - circle_size) / 2
        add_circle(s, cx, top + Inches(0.22), circle_size, fill=GREEN, text=str(i + 1))

        t = s.shapes.add_textbox(x + Inches(0.08), top + Inches(0.8), card_w - Inches(0.16), Inches(1.0))
        add_text(t, [(step_title, 11, True, BLACK)], align=PP_ALIGN.CENTER)

        d = s.shapes.add_textbox(x + Inches(0.1), top + Inches(1.9), card_w - Inches(0.2), Inches(1.55))
        add_text(d, [(desc, 10, False, GRAY)], align=PP_ALIGN.CENTER)

        sys_box = s.shapes.add_textbox(
            x + Inches(0.06), top + card_h - Inches(0.65), card_w - Inches(0.12), Inches(0.45)
        )
        add_text(sys_box, [(system, 11, True, GREEN)], align=PP_ALIGN.CENTER)

        add_appear_animation(s, card, i + 1)

    add_footer(s, prs)
    return s


def slide_gpc_workflow(prs):
    steps = [
        ("100% Ownership Check", "System verifies 100% ownership share in CLRMIS.", "CLRMIS"),
        ("Data Capture & English Name", "Auto-retrieves Urdu details; captures applicant English name.", "PZA App"),
        ("Mobile OTP Verification", "Validates 6-digit OTP sent to registered mobile.", "SMS Gateway"),
        ("Survey Representative", "Nominate optional representative for physical survey.", "PZA App"),
        ("SCO/ARC Verification & PSID", "Officer verifies record; generates PSID fee challan.", "LAMP"),
        ("Physical Survey & ADLR Approval", "Survey team completes survey; ADLR grants final approval.", "LAMP"),
        ("GPC Issued & Digital Vault", "Digital Certificate generated & filed in Digital Vault.", "Digital Vault"),
    ]
    slide_workflow(
        prs,
        "CORE APPLICATION MODULES",
        "Green Property Certificate (GPC) Workflow",
        "End-to-End Digital Certification Pipeline for 100% Owned Land",
        steps,
    )


def slide_hsms(prs):
    s = blank_slide(prs)
    add_eyebrow(s)
    add_title_block(
        s,
        "Approved Housing Society (HSMS)",
        "GIS satellite visualization of approved societies, block boundaries, and plot legal verification",
    )
    features = [
        ("GIS Satellite Interface", [
            "Interactive map with zoom and pan",
            "Full-screen polygon rendering",
            "Block and society boundary overlays",
        ]),
        ("Administrative Filters", [
            "Search by District and Tehsil",
            "Authority-Owned housing societies",
            "Private Housing Society filters",
        ]),
        ("Color-Coded Plot Statuses", [
            "Mortgage — Blue",
            "Legal Case — Orange",
            "Available — Green / Sold — Grey",
        ]),
        ("Block & Plot Cards", [
            "Plot size and key attributes",
            "Last transfer date",
            "Current legal / sale status",
        ]),
    ]
    positions = [
        (Inches(0.5), Inches(1.45)),
        (Inches(6.7), Inches(1.45)),
        (Inches(0.5), Inches(4.0)),
        (Inches(6.7), Inches(4.0)),
    ]
    for (x, y), (title, bullets) in zip(positions, features):
        add_bullet_card(s, x, y, Inches(5.9), Inches(2.25), title, bullets)
    add_footer(s, prs)


def slide_record_copy(prs):
    s = blank_slide(prs)
    add_eyebrow(s)
    add_title_block(
        s,
        "Get Record Copy Module",
        "Digitizing Fard, Registry & Mutation copies with integrated payments",
    )
    add_bullet_card(
        s,
        Inches(0.5),
        Inches(1.4),
        Inches(12.2),
        Inches(1.35),
        "Expanded Service Suite",
        [
            "Get Fard Copy  •  Get Registry Copy  •  Get Mutation (Intiqal) Copy  •  Crop Inspection Reports (Girdawari)",
            "Certified official copies with in-app fee settlement via Payment Gateway",
        ],
    )
    features = [
        ("Multi-Criteria Search", [
            "Lookup by e-Stamp Number",
            "Mutation Number / Registry Number",
            "Search by CNIC",
        ]),
        ("In-App Fee Settlement", [
            "Dedicated Payment Gateway",
            "Challan fee via Debit/Credit Card",
            "Instant receipt & vault storage",
        ]),
        ("Source Systems", [
            "CLRMIS integration for land records",
            "e-Stamp linkage where applicable",
            "Secure digital delivery to citizen",
        ]),
    ]
    for i, (title, bullets) in enumerate(features):
        x = Inches(0.5) + i * Inches(4.15)
        add_bullet_card(s, x, Inches(3.05), Inches(3.95), Inches(2.8), title, bullets)
    add_footer(s, prs)


def slide_payments(prs):
    s = blank_slide(prs)
    add_eyebrow(s)
    add_title_block(
        s,
        "Payment Gateway Module",
        "Online payment of challans for Fard, e-Stamp, copy fee, and more",
    )
    add_bullet_card(
        s,
        Inches(0.5),
        Inches(1.35),
        Inches(5.9),
        Inches(2.35),
        "AS-IS — Manual Payment Bottlenecks",
        [
            "Only basic Get Record Copy & static history",
            "No in-app payment gateway",
            "Citizens visited bank booths for challan fee",
            "Limited modules and fragmented experience",
        ],
    )
    add_bullet_card(
        s,
        Inches(6.7),
        Inches(1.35),
        Inches(5.9),
        Inches(2.35),
        "TO-BE — Integrated Payment Gateway",
        [
            "CLRMIS + e-Stamp integrated fee collection",
            "Debit/Credit card payments in-app & web",
            "Automated Digital Vault for receipts",
            "Future expansion to RCMS, HSMS, Arazi Muawin",
        ],
    )

    note = s.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.3))
    add_text(note, [("Three Integrated Components", 13, True, GREEN_DARK)])

    comps = [
        ("1. Citizen Web Portal", "Enter challan number, fetch details, and pay fee online"),
        ("2. Punjab Zameen App", "Payment card on home + Record Copy tab integration"),
        ("3. Admin Portal", "Monitor, verify, reconcile, and audit all transactions"),
    ]
    for i, (title, desc) in enumerate(comps):
        x = Inches(0.5) + i * Inches(4.15)
        add_rounded_rect(s, x, Inches(4.3), Inches(3.95), Inches(1.45), fill=GREEN_SOFT, line=GREEN_BORDER)
        t = s.shapes.add_textbox(x + Inches(0.2), Inches(4.45), Inches(3.55), Inches(0.35))
        add_text(t, [(title, 13, True, GREEN_DARK)])
        d = s.shapes.add_textbox(x + Inches(0.2), Inches(4.85), Inches(3.55), Inches(0.7))
        add_text(d, [(desc, 12, False, GRAY)])
    add_footer(s, prs)


def slide_payment_sources(prs):
    s = blank_slide(prs)
    add_eyebrow(s)
    add_title_block(
        s,
        "Payment Gateway — Supported Fees & Future Scope",
        "Current collections and planned expansion across PLRA systems",
    )
    add_bullet_card(
        s,
        Inches(0.5),
        Inches(1.4),
        Inches(6.0),
        Inches(3.35),
        "Currently Supported — CLRMIS",
        [
            "Fard Fee (Copy Fee of Fard)",
            "Mutation Fee",
            "Crop Inspection Fee",
            "Registry Fee",
        ],
    )
    add_bullet_card(
        s,
        Inches(6.8),
        Inches(1.4),
        Inches(5.9),
        Inches(3.35),
        "Currently Supported — E-Stamp",
        [
            "Fee of all E-Stamp instruments / papers",
            "Online settlement without bank booth visits",
            "Receipt available for download & vault",
            "Audit-ready transaction logs",
        ],
    )
    add_rounded_rect(s, Inches(0.5), Inches(5.0), Inches(12.2), Inches(1.15), fill=GREEN_SOFT, line=GREEN_BORDER)
    t = s.shapes.add_textbox(Inches(0.75), Inches(5.15), Inches(11.7), Inches(0.3))
    add_text(t, [("Future Payment Gateway Integrations", 13, True, GREEN_DARK)])
    d = s.shapes.add_textbox(Inches(0.75), Inches(5.55), Inches(11.7), Inches(0.4))
    add_text(
        d,
        [("RCMS  •  HSMS  •  Arazi Muawin  •  Additional PLRA systems via web portal & Punjab Zameen app", 13, False, GRAY)],
    )
    add_footer(s, prs)


def slide_payment_steps(prs):
    steps = [
        ("Enter Challan Number", "Citizen enters PSID / challan number on portal or app.", "Payment Gateway"),
        ("Fetch & Review Details", "System retrieves and displays challan amount & particulars.", "Payment Gateway"),
        ("Click Pay Now", "Proceed to the bank payment page securely.", "Bank Page"),
        ("Enter Card Details", "Provide Debit / Credit card information.", "Bank Page"),
        ("Click Pay", "Authorize and complete the payment transaction.", "Bank Page"),
        ("Download Receipt", "Obtain payment receipt; store in Digital Vault.", "Digital Vault"),
    ]
    slide_workflow(
        prs,
        "CORE APPLICATION MODULES",
        "Online Payment Steps",
        "Citizen journey for challan fee payment via web portal & Punjab Zameen",
        steps,
    )


def slide_cms(prs):
    s = blank_slide(prs)
    add_eyebrow(s)
    add_title_block(
        s,
        "Complaints Module (CMS)",
        "Citizen grievance redressal integrated with CMS Helpdesk workflows",
    )
    add_bullet_card(
        s,
        Inches(0.5),
        Inches(1.4),
        Inches(5.9),
        Inches(2.4),
        "Option 1 — Track My Complaint",
        [
            "Check complaint status in real time",
            "Statuses: Pending, Resolved, Rejected",
            "View timeline and official updates",
            "Push notifications for progress changes",
        ],
    )
    add_bullet_card(
        s,
        Inches(6.7),
        Inches(1.4),
        Inches(5.9),
        Inches(2.4),
        "Option 2 — Register a Complaint",
        [
            "Select complaint type / category",
            "Fill required workflow form fields",
            "Review → Edit (if needed) → Submit",
            "Receive CMS reference number on success",
        ],
    )
    add_bullet_card(
        s,
        Inches(0.5),
        Inches(4.05),
        Inches(12.2),
        Inches(1.75),
        "Citizen Portal Capabilities",
        [
            "Register / log in, file complaints, monitor progress, and communicate with authorities",
            "Synchronized category & sub-category taxonomy from CMS",
            "Attachments supported (documents/photos, up to 5MB) during complaint lodging",
        ],
    )
    add_footer(s, prs)


def slide_cms_workflow(prs):
    steps = [
        ("Category Selection", "Synchronized category & sub-category taxonomy from CMS.", "CMS"),
        ("Form & Attachment Upload", "Enter complaint details; attach documents/photos (up to 5MB).", "PZA App"),
        ("CMS Registration & Ref #", "Payload submitted via secure API; returns unique CMS Ref Number.", "CMS"),
        ("Helpdesk Assignment", "CMS routes complaint to designated officer or district desk.", "CMS Helpdesk"),
        ("Live Progress Tracking", "Real-time timeline updates delivered via push notifications.", "PZA App"),
        ("Resolution & Feedback", "Citizen views official resolution text; rates resolution satisfaction.", "CMS"),
    ]
    slide_workflow(
        prs,
        "CORE APPLICATION MODULES",
        "Complaints Module",
        "Citizen Grievance Redressal Integrated with CMS Helpdesk Workflows",
        steps,
    )


def slide_arazi(prs):
    s = blank_slide(prs)
    add_eyebrow(s)
    add_title_block(
        s,
        "Arazi Muawin Module",
        "Locate authorized franchise centers and submit new franchise applications",
    )
    features = [
        ("GPS Proximity Engine", [
            "Nearest Arazi Muawin office on interactive map",
            "GPS proximity enhancements planned",
            "Citizen-friendly discovery of centers",
        ]),
        ("Dual View Mode", [
            "Interactive GIS Map View",
            "Filtered List View",
            "Seamless toggle between modes",
        ]),
        ("Direct Device Actions", [
            "One-tap Call opens phone dialer",
            "Navigate launches Google Maps",
            "Faster real-world assistance",
        ]),
        ("Franchise Application Portal", [
            "Capture personal data & CNIC",
            "Office GPS & infrastructure details",
            "Tenancy and application particulars",
        ]),
    ]
    positions = [
        (Inches(0.5), Inches(1.45)),
        (Inches(6.7), Inches(1.45)),
        (Inches(0.5), Inches(4.0)),
        (Inches(6.7), Inches(4.0)),
    ]
    for (x, y), (title, bullets) in zip(positions, features):
        add_bullet_card(s, x, y, Inches(5.9), Inches(2.25), title, bullets)
    add_footer(s, prs)


def slide_rcms(prs):
    s = blank_slide(prs)
    add_eyebrow(s)
    add_title_block(
        s,
        "Revenue Court Cases (RCMS)",
        "Search cases, inspect cause lists, track hearings, and view judgments",
    )
    features = [
        ("Multi-Criteria Case Search", [
            "Search by Case ID",
            "Search by Case Title",
            "Search by Case Type",
        ]),
        ("Complete Party Profiles", [
            "Clear Petitioner Name breakdown",
            "Clear Respondent Name breakdown",
            "Party visibility for litigants",
        ]),
        ("Court Summary", [
            "Court Name & Presiding Judge",
            "Date of Institution",
            "Next Hearing Date",
        ]),
        ("Chronological Case Timeline", [
            "Visual progression of case events",
            "From Case Filed to Final Judgment",
            "Hearing schedule visibility",
        ]),
    ]
    positions = [
        (Inches(0.5), Inches(1.45)),
        (Inches(6.7), Inches(1.45)),
        (Inches(0.5), Inches(4.0)),
        (Inches(6.7), Inches(4.0)),
    ]
    for (x, y), (title, bullets) in zip(positions, features):
        add_bullet_card(s, x, y, Inches(5.9), Inches(2.25), title, bullets)
    add_footer(s, prs)


def slide_ereg(prs):
    s = blank_slide(prs)
    add_eyebrow(s)
    add_title_block(
        s,
        "e-Registration Module",
        "Property registration pipeline tracking & certified deed copies",
    )
    cards = [
        ("Dynamic Task Search", [
            "Search active registration tasks by CNIC",
            "Lookup by Task ID",
            "Lookup by e-Stamp Challan Number",
        ]),
        ("Registration Summary", [
            "Buyer Name and Seller Name",
            "Property Location",
            "Registration Type and Date",
        ]),
        ("Vault Integration", [
            "Certified registry deeds auto-saved",
            "Filed instantly in Digital Vault",
            "Downloadable signed documents",
        ]),
    ]
    for i, (title, bullets) in enumerate(cards):
        x = Inches(0.5) + i * Inches(4.15)
        add_bullet_card(s, x, Inches(1.55), Inches(3.95), Inches(4.15), title, bullets)
    add_footer(s, prs)


def slide_estamp(prs):
    s = blank_slide(prs)
    add_eyebrow(s)
    add_title_block(
        s,
        "e-Stamp Module",
        "Stamp instruments and online fee collection through Payment Gateway",
    )
    cards = [
        ("System Access", [
            "Dedicated e-Stamp card on home screen",
            "Access stamp instruments / papers",
            "Integrated with Payment Gateway",
        ]),
        ("Payment Linkage", [
            "Enter challan and pay online",
            "Debit / Credit card settlement",
            "Downloadable payment receipt",
        ]),
        ("Citizen Outcome", [
            "No physical bank booth for stamp fees",
            "Faster registration journeys",
            "Transparent fee tracking & audit trail",
        ]),
    ]
    for i, (title, bullets) in enumerate(cards):
        x = Inches(0.5) + i * Inches(4.15)
        add_bullet_card(s, x, Inches(1.55), Inches(3.95), Inches(4.15), title, bullets)
    add_footer(s, prs)


def slide_pts(prs):
    s = blank_slide(prs)
    add_eyebrow(s)
    add_title_block(
        s,
        "Property Tax Module (PTS)",
        "Excise integration for tax search, assessment verification & instant settlement",
    )
    note = s.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.3))
    add_text(note, [("Functionalities to be finalized — planned capabilities below", 12, True, GREEN)])

    features = [
        ("Multi-Criteria Tax Lookup", [
            "Search by Property Tax Number",
            "Search by Property ID",
            "Search by CNIC",
        ]),
        ("Assessment Details", [
            "Covered area & assessment year",
            "Current tax status",
            "Surcharges visibility",
        ]),
        ("Duplicate Payment Safeguard", [
            "Locks settled tax bills",
            "Prevents accidental re-payment",
            "Protects citizen transactions",
        ]),
        ("Tax Verification Certificate", [
            "Instant clear/not-clear status",
            "Current fiscal year verification",
            "Supports downstream land services",
        ]),
    ]
    positions = [
        (Inches(0.5), Inches(1.8)),
        (Inches(6.7), Inches(1.8)),
        (Inches(0.5), Inches(4.2)),
        (Inches(6.7), Inches(4.2)),
    ]
    for (x, y), (title, bullets) in zip(positions, features):
        add_bullet_card(s, x, y, Inches(5.9), Inches(2.05), title, bullets)
    add_footer(s, prs)


def slide_eleasing(prs):
    s = blank_slide(prs)
    add_eyebrow(s)
    add_title_block(
        s,
        "e-Leasing Module",
        "State land leasing schemes, security deposit settlement & live bidding",
    )
    features = [
        ("Existing Lessee Portal", [
            "View active agreements",
            "Lease durations & payment schedules",
            "Settle outstanding dues",
        ]),
        ("Scheme Search & Filter", [
            "Agricultural land schemes",
            "Commercial land schemes",
            "Industrial land schemes",
        ]),
        ("Security Deposit Gateway", [
            "Settle deposit via PSID or card",
            "Unlock bidding qualification",
            "Transparent fee handling",
        ]),
        ("Live Electronic Bidding", [
            "Real-time auction room",
            "Live bid ticks & highest bid",
            "Countdown to close",
        ]),
    ]
    positions = [
        (Inches(0.5), Inches(1.45)),
        (Inches(6.7), Inches(1.45)),
        (Inches(0.5), Inches(4.0)),
        (Inches(6.7), Inches(4.0)),
    ]
    for (x, y), (title, bullets) in zip(positions, features):
        add_bullet_card(s, x, y, Inches(5.9), Inches(2.25), title, bullets)
    add_footer(s, prs)


def slide_community(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "PLATFORM MODULES")
    add_title_block(
        s,
        "Community Module",
        "Collaborative knowledge sharing, verified practitioner answers & versioned guides",
    )
    features = [
        ("Real-Time Duplicate Prevention", [
            "Auto-searches existing discussions",
            "Triggered while typing question title",
            "Reduces repeated queries",
        ]),
        ("Verified Practitioner Authority", [
            "Only certified practitioners publish official answers",
            "Moderators control quality",
            "Trusted guidance for citizens",
        ]),
        ("Threaded Discussions", [
            "Original askers post follow-ups",
            "Conversation stays in same thread",
            "Clear Q&A continuity",
        ]),
        ("Version-Controlled Guides", [
            "Official procedural guides",
            "Fees, steps, and document lists",
            "Revision notes maintained",
        ]),
    ]
    positions = [
        (Inches(0.5), Inches(1.45)),
        (Inches(6.7), Inches(1.45)),
        (Inches(0.5), Inches(4.0)),
        (Inches(6.7), Inches(4.0)),
    ]
    for (x, y), (title, bullets) in zip(positions, features):
        add_bullet_card(s, x, y, Inches(5.9), Inches(2.25), title, bullets)
    add_footer(s, prs)


def slide_vault(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "PLATFORM MODULES")
    add_title_block(
        s,
        "Digital Vault Module",
        "Beyond basic service history — secure automated digital locker for issued documents",
    )
    features = [
        ("Evolution Beyond Service History", [
            "Replaces passive history logs",
            "Complete downloadable certified documents",
            "Legally valid issued artifacts",
        ]),
        ("Automatic Filing Engine", [
            "Documents from 12+ government systems",
            "Archived instantly on issuance",
            "No manual citizen upload needed",
        ]),
        ("Categorized Folders", [
            "GPC, Registries, Mutations",
            "Fard & Crop Reports",
            "Property Tax Receipts & more",
        ]),
        ("QR Seals & PDF Sharing", [
            "Download original signed PDFs",
            "Instant QR authenticity verification",
            "Secure share-ready outputs",
        ]),
    ]
    positions = [
        (Inches(0.5), Inches(1.45)),
        (Inches(6.7), Inches(1.45)),
        (Inches(0.5), Inches(4.0)),
        (Inches(6.7), Inches(4.0)),
    ]
    for (x, y), (title, bullets) in zip(positions, features):
        add_bullet_card(s, x, y, Inches(5.9), Inches(2.25), title, bullets)
    add_footer(s, prs)


def slide_notifications(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "PLATFORM MODULES")
    add_title_block(
        s,
        "Notifications & Messaging",
        "Multi-channel push, SMS OTP, and transactional messaging system",
    )
    cards = [
        ("Push Notifications", [
            "Real-time OS push alerts",
            "Status changes & hearing dates",
            "Live bidding alerts",
        ]),
        ("Payment Receipts", [
            "Instant financial confirmations",
            "PSID and transaction references",
            "Supports vault filing workflows",
        ]),
        ("Proclamation Alerts", [
            "Automated subscription notices",
            "Public notices on linked properties",
            "Keeps citizens informed proactively",
        ]),
    ]
    for i, (title, bullets) in enumerate(cards):
        x = Inches(0.5) + i * Inches(4.15)
        add_bullet_card(s, x, Inches(1.55), Inches(3.95), Inches(4.15), title, bullets)
    add_footer(s, prs)


def slide_closing(prs):
    s = blank_slide(prs)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.9), prs.slide_width, Inches(3.4))
    band.fill.solid()
    band.fill.fore_color.rgb = GREEN
    band.line.fill.background()

    t = s.shapes.add_textbox(Inches(0.8), Inches(2.45), Inches(11.5), Inches(0.7))
    add_text(t, [("Thank You", 42, True, WHITE)], align=PP_ALIGN.CENTER)

    d = s.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(11.5), Inches(0.5))
    add_text(
        d,
        [("Punjab Zameen — Empowering Citizens Through Digital Land Services", 16, False, RGBColor(0xD8, 0xF0, 0xE4))],
        align=PP_ALIGN.CENTER,
    )

    f = s.shapes.add_textbox(Inches(0.8), Inches(4.05), Inches(11.5), Inches(0.4))
    add_text(
        f,
        [("PUNJAB LAND RECORDS AUTHORITY (PLRA)  |  GOVT OF THE PUNJAB", 12, True, WHITE)],
        align=PP_ALIGN.CENTER,
    )


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_agenda(prs)
    slide_why_remodel(prs)
    slide_ecosystem(prs)
    slide_home_cards(prs)
    slide_my_properties(prs)
    slide_gpc_workflow(prs)
    slide_hsms(prs)
    slide_record_copy(prs)
    slide_payments(prs)
    slide_payment_sources(prs)
    slide_payment_steps(prs)
    slide_cms(prs)
    slide_cms_workflow(prs)
    slide_arazi(prs)
    slide_rcms(prs)
    slide_ereg(prs)
    slide_estamp(prs)
    slide_pts(prs)
    slide_eleasing(prs)
    slide_community(prs)
    slide_vault(prs)
    slide_notifications(prs)
    slide_closing(prs)

    apply_transitions(prs)

    out = "/workspace/Punjab_Zameen_Application.pptx"
    prs.save(out)
    print(f"Saved {out} with {len(prs.slides)} slides")
    return out


if __name__ == "__main__":
    build()
